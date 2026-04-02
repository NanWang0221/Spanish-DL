#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Created on Thu Aug 21 17:32:44 2025

@author: willchang
"""

import pandas as pd
import random
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches


# Load your CSV file (replace with your actual filename)
df = pd.read_csv("/Users/willchang/Dropbox/Research/UC Irvine/Research/research projects/Xin/NIH accent adaptation/stimuli recording list.csv")

# Create a list of all words with labels for constraint checking
words = []
for i, row in df.iterrows():
    if pd.notna(row['critical_voiced']):
        words.append(('cvd', i, row['critical_voiced']))
    if pd.notna(row['critical_voiceless']):
        words.append(('cvl', i, row['critical_voiceless']))
    if pd.notna(row['fillerA']):
        words.append(('fa', i, row['fillerA']))
    if pd.notna(row['fillerB']):
        words.append(('fb', i, row['fillerB']))
    if pd.notna(row['catch']):
        words.append(('catch', i, row['catch']))

# Repeat each word 5 times
words_repeated = words * 5
random.shuffle(words_repeated)

# Function to check constraints
def valid_block_assignment(blocks, word, block_num):
    label, idx, token = word
    # Ensure no cvd/cvl conflict
    if label in ['cvd','cvl']:
        conflict_labels = ['cvd','cvl']
    elif label in ['fa','fb']:
        conflict_labels = ['fa','fb']
    else:
        conflict_labels = []

    if conflict_labels:
        for w in blocks[block_num]:
            if w[1] == idx and w[0] in conflict_labels:
                return False
    return True

def assign_words_equal_blocks(words_repeated, num_blocks=10, block_size=41):
    success = False
    while not success:
        blocks = {i: [] for i in range(num_blocks)}
        success = True
        random.shuffle(words_repeated)  # shuffle order to help randomness
        
        for word in words_repeated:
            placed = False
            random_blocks = list(range(num_blocks))
            random.shuffle(random_blocks)
            for b in random_blocks:
                # Check constraints and size
                if valid_block_assignment(blocks, word, b) and len(blocks[b]) < block_size:
                    blocks[b].append(word)
                    placed = True
                    break
            if not placed:
                # failed assignment, restart
                success = False
                break
    return blocks

# Run assignment
blocks = assign_words_equal_blocks(words_repeated, num_blocks=10, block_size=41)

# Convert to DataFrame with each block as a column
block_df = pd.DataFrame({f'block_{i+1}': [w[2] for w in blocks[i]] for i in range(10)})
block_df.to_csv('/Users/willchang/Dropbox/Research/UC Irvine/Research/research projects/Xin/NIH accent adaptation/all_blocks_equal_size.csv', index=False)

print("Randomized blocks saved in one file: all_blocks_equal_size.csv")



# Present word lists in powerpoint slides
stimfile = "/Users/willchang/Dropbox/Research/UC Irvine/Research/research projects/Xin/NIH accent adaptation/all_blocks_equal_size.csv"
df_stim = pd.read_csv(stimfile, header=0)  # first row is column names
prs = Presentation()

# Slide dimensions
slide_width = prs.slide_width
slide_height = prs.slide_height

def add_centered_slide(text, font_size=72, bold=False, color=(0, 0, 0)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide

    # Textbox size
    textbox_width = Inches(8)
    textbox_height = Inches(2)

    # Center position
    left = (slide_width - textbox_width) / 2
    top = (slide_height - textbox_height) / 2

    textbox = slide.shapes.add_textbox(left, top, textbox_width, textbox_height)
    tf = textbox.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    p.alignment = PP_ALIGN.CENTER

    return slide

for col in df_stim.columns:
    # Column title slide in bold red
    add_centered_slide(col, font_size=60, bold=True, color=(139, 0, 0))

    # Word slides
    for word in df_stim[col]:
        if pd.notna(word):
            add_centered_slide(word, font_size=72)

prs.save("/Users/willchang/Dropbox/Research/UC Irvine/Research/research projects/Xin/NIH accent adaptation/stimuli_recording.pptx")

